"""Generates docs/presentation/Municipal-Law-Skill.pptx from content.py.

Run: python docs/presentation/build_pptx.py
Requires: pip install python-pptx (already in requirements-dev.txt)
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

from content import SLIDES

# Color palette - matches docs/demo.html's dark theme for a consistent
# visual identity across every artifact this project ships.
BG = RGBColor(0x0B, 0x0D, 0x12)
PANEL = RGBColor(0x12, 0x15, 0x1C)
BORDER = RGBColor(0x26, 0x2B, 0x36)
TEXT = RGBColor(0xE8, 0xEA, 0xED)
MUTED = RGBColor(0x9A, 0xA4, 0xB2)
ACCENT = RGBColor(0x5B, 0x9D, 0xFF)
ACCENT_2 = RGBColor(0x7E, 0xE0, 0xC3)
RED = RGBColor(0xFF, 0x6B, 0x6B)

FONT = "Segoe UI"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

OUT_PATH = Path(__file__).parent / "Municipal-Law-Skill.pptx"


def _blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    return slide


def _textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return box, tf


def _set_run(run, text, size, color, bold=False, italic=False, font=FONT):
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def _footer(slide, index, total):
    box, tf = _textbox(slide, Inches(0.5), Inches(7.05), Inches(10.8), Inches(0.35))
    p = tf.paragraphs[0]
    r = p.add_run()
    _set_run(r, "Municipal Law Skill for Autonomous Agents", 10, MUTED)

    box2, tf2 = _textbox(slide, Inches(11.8), Inches(7.05), Inches(1.0), Inches(0.35))
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    _set_run(r2, f"{index}/{total}", 10, MUTED)


def _title_bar(slide, title, subtitle=None):
    box, tf = _textbox(slide, Inches(0.7), Inches(0.45), Inches(11.9), Inches(1.1))
    p = tf.paragraphs[0]
    r = p.add_run()
    _set_run(r, title, 30, TEXT, bold=True)
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        _set_run(r2, subtitle, 14, MUTED, italic=True)
    # accent underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.35), Inches(2.2), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def render_title(slide, s):
    box, tf = _textbox(slide, Inches(1), Inches(2.3), Inches(11.3), Inches(1.2), MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    _set_run(r, s["title"], 48, TEXT, bold=True)

    box2, tf2 = _textbox(slide, Inches(1), Inches(3.35), Inches(11.3), Inches(0.8))
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    _set_run(r2, s["subtitle"], 26, ACCENT, bold=True)

    box3, tf3 = _textbox(slide, Inches(1.5), Inches(4.4), Inches(10.3), Inches(0.9))
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    _set_run(r3, s["tagline"], 16, MUTED, italic=True)

    box4, tf4 = _textbox(slide, Inches(1), Inches(6.3), Inches(11.3), Inches(0.5))
    p4 = tf4.paragraphs[0]
    p4.alignment = PP_ALIGN.CENTER
    r4 = p4.add_run()
    _set_run(r4, s["footer"], 13, MUTED)
    p5 = tf4.add_paragraph()
    p5.alignment = PP_ALIGN.CENTER
    r5 = p5.add_run()
    _set_run(r5, s["url"], 14, ACCENT_2, bold=True)


def render_bullets(slide, s, index, total):
    _title_bar(slide, s["title"])
    box, tf = _textbox(slide, Inches(1), Inches(1.9), Inches(11.3), Inches(4.6))
    for i, b in enumerate(s["bullets"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(18)
        r = p.add_run()
        _set_run(r, "▸  " + b, 20, TEXT)
    _footer(slide, index, total)


def render_quote(slide, s, index, total):
    _title_bar(slide, s["title"])
    # quote panel
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.9), Inches(11.3), Inches(1.4))
    panel.fill.solid()
    panel.fill.fore_color.rgb = PANEL
    panel.line.color.rgb = ACCENT
    panel.line.width = Pt(1.5)
    tf = panel.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run()
    _set_run(r, "“" + s["quote"].strip("“”") + "”", 18, ACCENT_2, italic=True, bold=True)

    top = Inches(3.5)
    if s.get("flows"):
        for flow in s["flows"]:
            label_box, ltf = _textbox(slide, Inches(1), top, Inches(11.3), Inches(0.3))
            lp = ltf.paragraphs[0]
            lr = lp.add_run()
            _set_run(lr, flow["label"], 13, MUTED, bold=True)
            top = Emu(top + Inches(0.32))

            steps = flow["steps"]
            n = len(steps)
            box_w = Inches(11.3 / n * 0.82)
            gap = Inches(11.3 / n * 0.18)
            left = Inches(1)
            box_h = Inches(0.55)
            muted = flow.get("muted")
            for i, step in enumerate(steps):
                shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, box_w, box_h)
                shp.fill.solid()
                shp.fill.fore_color.rgb = PANEL
                shp.line.color.rgb = BORDER if muted else ACCENT_2
                shp.line.width = Pt(1)
                stf = shp.text_frame
                stf.word_wrap = True
                stf.vertical_anchor = MSO_ANCHOR.MIDDLE
                sp = stf.paragraphs[0]
                sp.alignment = PP_ALIGN.CENTER
                sr = sp.add_run()
                _set_run(sr, step, 10, MUTED if muted else ACCENT_2, bold=not muted)
                if i < n - 1:
                    arrow_left = Emu(left + box_w)
                    arrow, atf = _textbox(slide, arrow_left, top, gap, box_h, MSO_ANCHOR.MIDDLE)
                    ap = atf.paragraphs[0]
                    ap.alignment = PP_ALIGN.CENTER
                    ar = ap.add_run()
                    _set_run(ar, "→", 14, MUTED)
                left = Emu(left + box_w + gap)
            top = Emu(top + box_h + Inches(0.15))
    else:
        box, tf2 = _textbox(slide, Inches(1), top, Inches(11.3), Inches(1.6))
        for i, b in enumerate(s.get("bullets", [])):
            p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            p2.space_after = Pt(14)
            r2 = p2.add_run()
            _set_run(r2, "▸  " + b, 18, TEXT)
        top = Emu(top + Inches(1.6))

    if s.get("emphasis"):
        emp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Emu(top + Inches(0.1)), Inches(11.3), Inches(0.9))
        emp.fill.solid()
        emp.fill.fore_color.rgb = ACCENT
        emp.line.fill.background()
        etf = emp.text_frame
        etf.word_wrap = True
        etf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ep = etf.paragraphs[0]
        ep.alignment = PP_ALIGN.CENTER
        er = ep.add_run()
        _set_run(er, s["emphasis"], 20, RGBColor(0x0B, 0x0D, 0x12), bold=True)

    _footer(slide, index, total)


def render_myth(slide, s, index, total):
    _title_bar(slide, s["title"])

    panel_w = Inches(5.15)
    panel_h = Inches(2.0)
    top = Inches(1.9)

    wrong = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), top, panel_w, panel_h)
    wrong.fill.solid()
    wrong.fill.fore_color.rgb = PANEL
    wrong.line.color.rgb = RED
    wrong.line.width = Pt(1.5)
    wtf = wrong.text_frame
    wtf.word_wrap = True
    wtf.vertical_anchor = MSO_ANCHOR.MIDDLE
    wp = wtf.paragraphs[0]
    wp.alignment = PP_ALIGN.CENTER
    wr = wp.add_run()
    _set_run(wr, "✗", 26, RED, bold=True)
    wp2 = wtf.add_paragraph()
    wp2.alignment = PP_ALIGN.CENTER
    wr2 = wp2.add_run()
    _set_run(wr2, s["wrong_label"], 13, MUTED)
    wp3 = wtf.add_paragraph()
    wp3.alignment = PP_ALIGN.CENTER
    wr3 = wp3.add_run()
    _set_run(wr3, s["wrong_claim"], 20, RED, bold=True)

    arrow, atf = _textbox(slide, Inches(5.95), top, Inches(0.6), panel_h, MSO_ANCHOR.MIDDLE)
    ap = atf.paragraphs[0]
    ap.alignment = PP_ALIGN.CENTER
    ar = ap.add_run()
    _set_run(ar, "→", 26, ACCENT, bold=True)

    right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.65), top, panel_w, panel_h)
    right.fill.solid()
    right.fill.fore_color.rgb = PANEL
    right.line.color.rgb = ACCENT_2
    right.line.width = Pt(1.5)
    rtf = right.text_frame
    rtf.word_wrap = True
    rtf.vertical_anchor = MSO_ANCHOR.MIDDLE
    rp = rtf.paragraphs[0]
    rp.alignment = PP_ALIGN.CENTER
    rr = rp.add_run()
    _set_run(rr, "✓", 26, ACCENT_2, bold=True)
    rp2 = rtf.add_paragraph()
    rp2.alignment = PP_ALIGN.CENTER
    rr2 = rp2.add_run()
    _set_run(rr2, s["right_label"], 13, MUTED)
    rp3 = rtf.add_paragraph()
    rp3.alignment = PP_ALIGN.CENTER
    rr3 = rp3.add_run()
    _set_run(rr3, s["right_claim"], 20, ACCENT_2, bold=True)
    rp4 = rtf.add_paragraph()
    rp4.alignment = PP_ALIGN.CENTER
    rr4 = rp4.add_run()
    _set_run(rr4, s["right_detail"], 12, MUTED)

    box, tf = _textbox(slide, Inches(1), Inches(4.2), Inches(11.3), Inches(1.2))
    p = tf.paragraphs[0]
    r = p.add_run()
    _set_run(r, s["footer_line"], 15, MUTED)
    _footer(slide, index, total)


def render_comparison(slide, s, index, total):
    _title_bar(slide, s["title"])
    col_w = Inches(5.4)
    top = Inches(1.9)
    gap = Inches(0.5)

    left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), top, col_w, Inches(3.6))
    left.fill.solid()
    left.fill.fore_color.rgb = PANEL
    left.line.color.rgb = BORDER
    ltf = left.text_frame
    ltf.word_wrap = True
    ltf.margin_left = Inches(0.3)
    ltf.margin_top = Inches(0.25)
    lp = ltf.paragraphs[0]
    lp.alignment = PP_ALIGN.CENTER
    lr = lp.add_run()
    _set_run(lr, s["left_label"], 17, MUTED, bold=True)
    for item in s["left_items"]:
        lp2 = ltf.add_paragraph()
        lp2.space_before = Pt(12)
        lr2 = lp2.add_run()
        _set_run(lr2, "✗  " + item, 16, TEXT)

    right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(Inches(0.9) + col_w + gap), top, col_w, Inches(3.6))
    right.fill.solid()
    right.fill.fore_color.rgb = PANEL
    right.line.color.rgb = ACCENT_2
    rtf = right.text_frame
    rtf.word_wrap = True
    rtf.margin_left = Inches(0.3)
    rtf.margin_top = Inches(0.25)
    rp = rtf.paragraphs[0]
    rp.alignment = PP_ALIGN.CENTER
    rr = rp.add_run()
    _set_run(rr, s["right_label"], 17, ACCENT_2, bold=True)
    for item in s["right_items"]:
        rp2 = rtf.add_paragraph()
        rp2.space_before = Pt(12)
        rr2 = rp2.add_run()
        _set_run(rr2, "✓  " + item, 16, TEXT)

    _footer(slide, index, total)


def render_diagram(slide, s, index, total):
    _title_bar(slide, s["title"])
    if s.get("lead"):
        lead_box, ltf = _textbox(slide, Inches(1), Inches(1.75), Inches(11.3), Inches(0.6))
        lp = ltf.paragraphs[0]
        lp.alignment = PP_ALIGN.CENTER
        lr = lp.add_run()
        _set_run(lr, s["lead"], 20, ACCENT_2, bold=True)
        diagram_top = 2.9
    else:
        diagram_top = 2.6
    steps = s["diagram"]
    n = len(steps)
    box_w = Inches(11.3 / n * 0.82)
    gap = Inches(11.3 / n * 0.18)
    left = Inches(1)
    top = Inches(diagram_top)
    for i, step in enumerate(steps):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, box_w, Inches(1.1))
        shp.fill.solid()
        shp.fill.fore_color.rgb = PANEL
        shp.line.color.rgb = ACCENT if i in (0, n - 1) else BORDER
        shp.line.width = Pt(1.25)
        tf = shp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        _set_run(r, step, 12, TEXT, bold=(i in (0, n - 1)))
        if i < n - 1:
            arrow_left = Emu(left + box_w)
            arrow, atf = _textbox(slide, arrow_left, Emu(top + Inches(0.25)), gap, Inches(0.6), MSO_ANCHOR.MIDDLE)
            ap = atf.paragraphs[0]
            ap.alignment = PP_ALIGN.CENTER
            ar = ap.add_run()
            _set_run(ar, "→", 22, ACCENT, bold=True)
        left = Emu(left + box_w + gap)

    box2, tf2 = _textbox(slide, Inches(1), Emu(top + Inches(1.7)), Inches(11.3), Inches(1.2), MSO_ANCHOR.MIDDLE)
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    _set_run(r2, s["callout"], 18, ACCENT_2, italic=True, bold=True)
    _footer(slide, index, total)


def render_stats(slide, s, index, total):
    _title_bar(slide, s["title"])
    groups = s["stat_groups"]
    n = len(groups)
    col_w = Inches(11.3 / n - 0.3)
    left = Inches(1)
    for group in groups:
        panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(2.0), col_w, Inches(2.8))
        panel.fill.solid()
        panel.fill.fore_color.rgb = PANEL
        panel.line.color.rgb = BORDER
        tf = panel.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        _set_run(r, group["label"], 16, MUTED, bold=True)
        for num, label in group["stats"]:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(10)
            r2 = p2.add_run()
            _set_run(r2, num, 34, ACCENT, bold=True)
            p3 = tf.add_paragraph()
            p3.alignment = PP_ALIGN.CENTER
            r3 = p3.add_run()
            _set_run(r3, label, 14, TEXT)
        left = Emu(left + col_w + Inches(0.3))

    box, tf2 = _textbox(slide, Inches(1), Inches(5.2), Inches(11.3), Inches(0.6), MSO_ANCHOR.MIDDLE)
    p = tf2.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    _set_run(r, s["footnote"], 16, ACCENT_2, italic=True)

    if s.get("callout"):
        box_c, tf_c = _textbox(slide, Inches(1), Inches(5.85), Inches(11.3), Inches(0.6), MSO_ANCHOR.MIDDLE)
        pc = tf_c.paragraphs[0]
        pc.alignment = PP_ALIGN.CENTER
        rc = pc.add_run()
        _set_run(rc, s["callout"], 18, TEXT, bold=True)

    _footer(slide, index, total)


def render_capability_table(slide, s, index, total):
    _title_bar(slide, s["title"], s.get("subtitle"))
    rows = s["rows"]
    top = Inches(2.0) if not s.get("subtitle") else Inches(2.2)
    table_shape = slide.shapes.add_table(len(rows) + 1, 3, Inches(0.7), top, Inches(12.0), Inches(3.8))
    table = table_shape.table
    table.columns[0].width = Inches(2.6)
    table.columns[1].width = Inches(2.6)
    table.columns[2].width = Inches(6.8)

    headers = ["Capability", "Endpoint", "Returns"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        tf = cell.text_frame
        tf.margin_top = Pt(4)
        tf.margin_bottom = Pt(4)
        p = tf.paragraphs[0]
        r = p.add_run()
        _set_run(r, h, 14, RGBColor(0x0B, 0x0D, 0x12), bold=True)

    for i, (capability, endpoint, returns) in enumerate(rows, start=1):
        for c, val in enumerate([capability, endpoint, returns]):
            cell = table.cell(i, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL if i % 2 else BG
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_top = Pt(6)
            tf.margin_bottom = Pt(6)
            p = tf.paragraphs[0]
            r = p.add_run()
            color = ACCENT_2 if c == 0 else (TEXT if c == 1 else MUTED)
            _set_run(r, val, 13 if c != 2 else 12.5, color, bold=(c in (0, 1)))

    if s.get("admin_note"):
        box, tf2 = _textbox(slide, Inches(0.7), Emu(top + Inches(3.9)), Inches(12.0), Inches(0.8))
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        _set_run(r2, s["admin_note"], 10, MUTED, italic=True)

    _footer(slide, index, total)


def render_story(slide, s, index, total):
    _title_bar(slide, s["title"])
    box, tf = _textbox(slide, Inches(1), Inches(1.75), Inches(11.3), Inches(0.7))
    p = tf.paragraphs[0]
    r = p.add_run()
    _set_run(r, s["question"], 22, ACCENT_2, italic=True, bold=True)

    flow = s["flow"]
    citation_index = s.get("citation_index")
    n = len(flow)
    box_w = Inches(11.3 / n * 0.82)
    gap = Inches(11.3 / n * 0.18)
    left = Inches(1)
    top = Inches(2.9)
    box_h = Inches(1.3)
    for i, step in enumerate(flow):
        is_cit = i == citation_index
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, box_w, box_h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = PANEL
        shp.line.color.rgb = ACCENT_2 if is_cit else BORDER
        shp.line.width = Pt(2 if is_cit else 1.25)
        tf2 = shp.text_frame
        tf2.word_wrap = True
        tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf2.margin_left = Inches(0.1)
        tf2.margin_right = Inches(0.1)
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        _set_run(r2, step, 13 if is_cit else 12, ACCENT_2 if is_cit else TEXT, bold=is_cit)
        if i < n - 1:
            arrow_left = Emu(left + box_w)
            arrow, atf = _textbox(slide, arrow_left, top, gap, box_h, MSO_ANCHOR.MIDDLE)
            ap = atf.paragraphs[0]
            ap.alignment = PP_ALIGN.CENTER
            ar = ap.add_run()
            _set_run(ar, "→", 20, ACCENT, bold=True)
        left = Emu(left + box_w + gap)

    if s.get("answer"):
        box3, tf3 = _textbox(slide, Inches(1.2), Emu(top + box_h + Inches(0.3)), Inches(10.9), Inches(1.0), MSO_ANCHOR.MIDDLE)
        p3 = tf3.paragraphs[0]
        p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run()
        _set_run(r3, s["answer"], 17, TEXT, italic=True)

    _footer(slide, index, total)


def render_closing(slide, s, index, total):
    box, tf = _textbox(slide, Inches(1), Inches(1.6), Inches(11.3), Inches(1.2), MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    _set_run(r, s["title"], 34, TEXT, bold=True)

    box2, tf2 = _textbox(slide, Inches(2), Inches(3.0), Inches(9.3), Inches(2.0))
    for i, b in enumerate(s["bullets"]):
        p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_after = Pt(10)
        r2 = p2.add_run()
        _set_run(r2, b, 18, ACCENT_2)

    box3, tf3 = _textbox(slide, Inches(1), Inches(5.4), Inches(11.3), Inches(1.0), MSO_ANCHOR.MIDDLE)
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    _set_run(r3, s["url"], 20, ACCENT, bold=True)
    p4 = tf3.add_paragraph()
    p4.alignment = PP_ALIGN.CENTER
    r4 = p4.add_run()
    _set_run(r4, s["repo"], 14, MUTED)
    _footer(slide, index, total)


RENDERERS = {
    "title": lambda slide, s, i, n: render_title(slide, s),
    "bullets": render_bullets,
    "quote": render_quote,
    "myth": render_myth,
    "comparison": render_comparison,
    "diagram": render_diagram,
    "stats": render_stats,
    "capability_table": render_capability_table,
    "story": render_story,
    "closing": render_closing,
}


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = len(SLIDES)
    for i, s in enumerate(SLIDES, start=1):
        slide = _blank_slide(prs)
        RENDERERS[s["kind"]](slide, s, i, total)

    prs.save(OUT_PATH)
    print(f"wrote {OUT_PATH} ({total} slides)")


if __name__ == "__main__":
    build()
